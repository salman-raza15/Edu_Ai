"""
loaders.py
-----------
This module detects and parses input files (CSV, PDF, ZIP), and extracts
data from all of them into a single UNIFIED format.

Unified record format (every row must follow this shape):
{
    "student_id": str,
    "student_name": str,
    "course_id": str,
    "cohort_id": str,
    "assignment_id": str,
    "criteria": str,
    "max_marks": float,
    "obtained_marks": float,
    "evaluation_date": str (YYYY-MM-DD)
}
"""

import os
import zipfile
import io
import re
import pandas as pd
import pdfplumber
from docx import Document
from pptx import Presentation

from .schema_mapper import map_columns_to_schema, apply_mapping

REQUIRED_COLUMNS = [
    "student_id", "student_name", "course_id", "cohort_id",
    "assignment_id", "criteria", "max_marks", "obtained_marks", "evaluation_date"
]


def detect_file_type(filename):
    """Detects file type based on its extension."""
    ext = os.path.splitext(filename)[1].lower()
    if ext == ".csv":
        return "csv"
    elif ext == ".pdf":
        return "pdf"
    elif ext == ".xlsx":
        return "excel"
    elif ext == ".docx":
        return "docx"
    elif ext == ".pptx":
        return "pptx"
    elif ext == ".zip":
        return "zip"
    else:
        raise ValueError(
            f"Unsupported file type: {ext}. Supported: .csv, .pdf, .xlsx, .docx, .pptx, .zip"
        )


def parse_csv(file_obj_or_path):
    """
    Parses a CSV using pandas. Instead of requiring exact column names,
    this uses schema_mapper to figure out (via AI, or alias fallback)
    which actual column corresponds to each required field — so any
    instructor's CSV, with any column naming, can be used.
    """
    df = pd.read_csv(file_obj_or_path)
    mapping = map_columns_to_schema(df)
    return apply_mapping(df, mapping)


def _find_field(pattern, text, default=""):
    """Regex helper shared by every document-style extractor (PDF/DOCX/PPTX)."""
    m = re.search(pattern, text, re.IGNORECASE)
    return m.group(1).strip() if m else default


def _extract_header_fields(full_text):
    """
    Pulls the 'Key: Value' header fields (Student Name, Student ID, Course,
    Cohort, Assignment, Evaluation Date) out of free text. Shared by the
    PDF, DOCX and PPTX extractors since they all assume the same basic
    report layout.
    """
    return {
        "student_name": _find_field(r"Student Name:\s*(.+)", full_text),
        "student_id": _find_field(r"Student ID:\s*(\S+)", full_text),
        "course_id": _find_field(r"Course:\s*(\S+)", full_text),
        "cohort_id": _find_field(r"Cohort:\s*(\S+)", full_text),
        "assignment_id": _find_field(r"Assignment:\s*(\S+)", full_text),
        "evaluation_date": _find_field(r"Evaluation Date:\s*(\S+)", full_text),
    }


def _rows_from_tables(tables, header_fields):
    """
    Scans a list of table grids (each table = list of rows = list of cell
    strings) for one whose header row contains Criteria / Max Marks /
    Obtained Marks columns, and turns the data rows into unified records.
    Shared by the PDF, DOCX and PPTX extractors.
    """
    rows = []
    for table in tables:
        if not table or len(table) < 2:
            continue
        header = [str(h or "").strip().lower() for h in table[0]]
        if "criteria" in header and any("max" in h for h in header) and any("obtained" in h for h in header):
            for row in table[1:]:
                if not row or len(row) < 3:
                    continue
                criteria, max_m, obt_m = row[0], row[1], row[2]
                try:
                    rows.append({
                        **header_fields,
                        "criteria": str(criteria).strip(),
                        "max_marks": float(max_m),
                        "obtained_marks": float(obt_m),
                    })
                except (ValueError, TypeError):
                    continue  # Skip unparsable rows instead of crashing
    return rows


def extract_pdf(file_obj_or_path):
    """
    Extracts text and table data from a PDF and converts it into unified rows.
    NOTE: This is a BASIC extractor that assumes the PDF contains:
      - Header lines following a 'Key: Value' pattern (Student Name, Student ID, etc.)
      - A table with Criteria / Max Marks / Obtained Marks columns

    Supports single-student PDFs AND multi-student "batch" PDFs — any number
    of student blocks (header + table) per page, in any layout: one per
    page, several per page, or spanning pages. Each table is matched to the
    header text that appears directly above it on the page (using each
    table's vertical position), not just "the first header found in the
    file" — so student records are never mixed up together regardless of
    how many appear on the same page.

    If a table has no header text of its own above it (e.g. a table that
    continues onto the next page), the most recently seen header is carried
    forward.
    """
    rows = []
    last_header = None
    with pdfplumber.open(file_obj_or_path) as pdf:
        for page in pdf.pages:
            found_tables = page.find_tables()
            if not found_tables:
                continue
            found_tables.sort(key=lambda t: t.bbox[1])  # top-to-bottom order

            prev_bottom = 0
            for t in found_tables:
                top = max(t.bbox[1], prev_bottom)
                if top > prev_bottom:
                    crop = page.within_bbox((0, prev_bottom, page.width, top))
                    header_text = crop.extract_text() or ""
                else:
                    header_text = ""
                page_header = _extract_header_fields(header_text)

                if page_header.get("student_id"):
                    header_fields = page_header
                    last_header = page_header
                elif last_header:
                    header_fields = last_header
                else:
                    header_fields = page_header  # no student_id seen yet at all

                rows.extend(_rows_from_tables([t.extract()], header_fields))
                prev_bottom = t.bbox[3]

    if not rows:
        raise ValueError("No valid data could be extracted from the PDF. "
                          "The format may differ from what was expected.")

    return pd.DataFrame(rows)


def parse_excel(file_obj_or_path):
    """
    Parses an Excel file (.xlsx). Every sheet is read and mapped to
    the schema independently (via schema_mapper, same as CSV) since some
    instructors keep one sheet per course/cohort, then all sheets are
    combined into one unified DataFrame.
    """
    sheets = pd.read_excel(file_obj_or_path, sheet_name=None)
    all_dfs = []
    for sheet_name, sheet_df in sheets.items():
        sheet_df = sheet_df.dropna(how="all")
        if sheet_df.empty:
            continue
        mapping = map_columns_to_schema(sheet_df)
        all_dfs.append(apply_mapping(sheet_df, mapping))

    if not all_dfs:
        raise ValueError("No data could be found in the Excel file.")

    return pd.concat(all_dfs, ignore_index=True)


def extract_docx(file_obj_or_path):
    """
    Extracts data from a Word (.docx) report.
    NOTE: Same BASIC assumption as extract_pdf():
      - 'Key: Value' header lines (Student Name, Student ID, etc.) in the
        document's paragraph text
      - A table with Criteria / Max Marks / Obtained Marks columns
    """
    doc = Document(file_obj_or_path)
    full_text = "\n".join(p.text for p in doc.paragraphs)
    header_fields = _extract_header_fields(full_text)

    tables = [
        [[cell.text.strip() for cell in row.cells] for row in table.rows]
        for table in doc.tables
    ]
    rows = _rows_from_tables(tables, header_fields)

    if not rows:
        raise ValueError("No valid data could be extracted from the DOCX file. "
                          "The format may differ from what was expected.")

    return pd.DataFrame(rows)


def extract_pptx(file_obj_or_path):
    """
    Extracts data from a PowerPoint (.pptx) report.
    NOTE: Same BASIC assumption as extract_pdf():
      - 'Key: Value' header lines (Student Name, Student ID, etc.) in any
        text box on any slide
      - A table (on any slide) with Criteria / Max Marks / Obtained Marks
        columns
    """
    prs = Presentation(file_obj_or_path)

    full_text_parts = []
    tables = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                full_text_parts.append(shape.text_frame.text)
            if shape.has_table:
                grid = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                tables.append(grid)

    header_fields = _extract_header_fields("\n".join(full_text_parts))
    rows = _rows_from_tables(tables, header_fields)

    if not rows:
        raise ValueError("No valid data could be extracted from the PPTX file. "
                          "The format may differ from what was expected.")

    return pd.DataFrame(rows)


_EXTRACTORS = {
    "csv": parse_csv,
    "pdf": extract_pdf,
    "excel": parse_excel,
    "docx": extract_docx,
    "pptx": extract_pptx,
}


def parse_zip(file_path):
    """Finds any supported file (CSV/PDF/XLSX/DOCX/PPTX) inside a ZIP and combines them into one unified DataFrame."""
    all_dfs = []
    with zipfile.ZipFile(file_path, "r") as z:
        for name in z.namelist():
            if name.endswith("/") or name.startswith("__MACOSX"):
                continue
            try:
                ftype = detect_file_type(name)
            except ValueError:
                continue  # Unknown file type, skip
            if ftype == "zip":
                continue  # Skip nested zips

            with z.open(name) as f:
                content = io.BytesIO(f.read())
                all_dfs.append(_EXTRACTORS[ftype](content))

    if not all_dfs:
        raise ValueError("No valid CSV, PDF, XLSX, DOCX or PPTX files found inside the ZIP.")

    return pd.concat(all_dfs, ignore_index=True)


def load_data(file_path):
    """
    MAIN ENTRY POINT — pass in any supported file (CSV/PDF/XLSX/DOCX/PPTX/ZIP)
    and this will return a unified pandas DataFrame.

    `file_path` can be a real path, or a file-like object (e.g. BytesIO) —
    in that case pass the type explicitly since there's no extension to
    detect from, or wrap it via load_data_from_stream() instead.
    """
    ftype = detect_file_type(file_path)

    if ftype == "zip":
        df = parse_zip(file_path)
    else:
        df = _EXTRACTORS[ftype](file_path)

    df, warnings = clean_and_validate(df)

    if warnings:
        print(f"\n[DATA QUALITY WARNING] {len(warnings)} row(s) had issues and were handled:")
        for w in warnings:
            print(f"  - {w}")
        print()

    return df


def load_data_from_stream(filename, file_obj):
    """
    Same as load_data(), but for use when the file arrives as an in-memory
    stream with a separate original filename (e.g. a FastAPI UploadFile) —
    there's no real path on disk to detect the extension from, so the
    filename is passed separately.
    """
    ftype = detect_file_type(filename)

    if ftype == "zip":
        # zipfile needs a seekable file-like object; BytesIO already is one
        df = parse_zip(file_obj)
    else:
        df = _EXTRACTORS[ftype](file_obj)

    df, warnings = clean_and_validate(df)

    if warnings:
        print(f"\n[DATA QUALITY WARNING] {len(warnings)} row(s) had issues and were handled:")
        for w in warnings:
            print(f"  - {w}")
        print()

    return df, warnings


def clean_and_validate(df: pd.DataFrame):
    """
    Cleans the data and checks for common real-world data issues.
    Returns (cleaned_df, list_of_warning_messages) so problems are visible
    instead of silently producing wrong report numbers.
    """
    warnings = []
    original_len = len(df)

    # 1. Convert marks to numeric; anything non-numeric becomes NaN
    df["max_marks"] = pd.to_numeric(df["max_marks"], errors="coerce")
    df["obtained_marks"] = pd.to_numeric(df["obtained_marks"], errors="coerce")

    # 2. Drop rows with missing/non-numeric marks
    missing_mask = df["max_marks"].isna() | df["obtained_marks"].isna()
    if missing_mask.sum() > 0:
        warnings.append(
            f"{missing_mask.sum()} row(s) dropped due to missing or invalid marks "
            f"(student_id(s): {df.loc[missing_mask, 'student_id'].tolist()})"
        )
    df = df[~missing_mask]

    # 3. Drop/flag rows where obtained_marks > max_marks (impossible — data entry error)
    invalid_score_mask = df["obtained_marks"] > df["max_marks"]
    if invalid_score_mask.sum() > 0:
        warnings.append(
            f"{invalid_score_mask.sum()} row(s) dropped because obtained_marks > max_marks "
            f"(student_id(s): {df.loc[invalid_score_mask, 'student_id'].tolist()}) — "
            f"likely a data entry error, needs manual review"
        )
    df = df[~invalid_score_mask]

    # 4. Standardize evaluation_date format (expects YYYY-MM-DD); flag rows that don't match
    date_pattern = r"^\d{4}-\d{2}-\d{2}$"
    bad_date_mask = ~df["evaluation_date"].astype(str).str.match(date_pattern)
    if bad_date_mask.sum() > 0:
        warnings.append(
            f"{bad_date_mask.sum()} row(s) dropped due to malformed evaluation_date "
            f"(expected YYYY-MM-DD) (student_id(s): {df.loc[bad_date_mask, 'student_id'].tolist()})"
        )
    df = df[~bad_date_mask]

    # 5. Check for exact duplicate rows
    dup_mask = df.duplicated(
        subset=["student_id", "assignment_id", "criteria", "evaluation_date"], keep="first"
    )
    if dup_mask.sum() > 0:
        warnings.append(
            f"{dup_mask.sum()} duplicate row(s) removed "
            f"(student_id(s): {df.loc[dup_mask, 'student_id'].tolist()})"
        )
    df = df[~dup_mask]

    if len(df) < original_len:
        warnings.append(f"Total: {original_len - len(df)} of {original_len} rows removed during cleaning.")

    return df, warnings