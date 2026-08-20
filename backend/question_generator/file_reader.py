import fitz
from docx import Document
from pptx import Presentation



def extract_text_from_pdf(file_path):

    text = ""

    pdf = fitz.open(file_path)

    for page in pdf:
        text += page.get_text()

    pdf.close()

    return text



def extract_text_from_docx(file_path):

    doc = Document(file_path)

    text = ""

    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"

    return text



def extract_text_from_pptx(file_path):

    presentation = Presentation(file_path)

    text = ""

    for slide in presentation.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                text += shape.text + "\n"

    return text



def extract_text(file_path):

    file_path = file_path.lower()


    if file_path.endswith(".pdf"):

        return extract_text_from_pdf(
            file_path
        )


    elif file_path.endswith(".docx"):

        return extract_text_from_docx(
            file_path
        )


    elif file_path.endswith(".pptx"):

        return extract_text_from_pptx(
            file_path
        )


    else:

        raise ValueError(
            "Unsupported file format"
        )